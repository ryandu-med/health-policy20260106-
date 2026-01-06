import streamlit as st
import os
import pandas as pd
from datetime import datetime
import pdfplumber
import docx
import openpyxl
from pptx import Presentation
import time
import plotly.express as px
import re
import logging
import base64
import pickle

# --- 0. 基础设置与屏蔽警告 ---
logging.getLogger("pdfminer").setLevel(logging.ERROR)

st.set_page_config(
    page_title="健康融入所有政策知识库平台",
    layout="wide",
    page_icon="🏛️",
    initial_sidebar_state="expanded"
)


# --- 1. 全局样式注入 (保持政务风格不变) ---
def inject_custom_css():
    st.markdown("""
        <style>
        .stApp { background-color: #f8f9fa; font-family: "Microsoft YaHei", "SimHei", sans-serif; }
        header[data-testid="stHeader"] { background-color: #1e50a2; }
        section[data-testid="stSidebar"] { background-color: #e9ecef; border-right: 1px solid #dcdcdc; }
        h1, h2, h3 { color: #1e50a2 !important; font-weight: 600; }
        .stButton > button { background-color: #1e50a2; color: white; border-radius: 2px; border: none; padding: 0.4rem 1rem; }
        .stButton > button:hover { background-color: #163e7f; color: white; }
        .file-row { background-color: white; padding: 10px; margin-bottom: 5px; border: 1px solid #e0e0e0; border-left: 4px solid #1e50a2; display: flex; align-items: center; }
        .breadcrumb { font-size: 14px; color: #666; padding: 10px 0; border-bottom: 2px solid #1e50a2; margin-bottom: 20px; }
        .block-container { padding-top: 2rem; }

        /* 目录样式微调 */
        .dir-header-3 { color: #1e50a2; font-size: 15px; font-weight: bold; margin: 15px 0 5px 0; padding-left: 10px; border-left: 3px solid #b22222; }
        .dir-header-4 { color: #333; font-size: 14px; font-weight: bold; margin: 10px 0 5px 20px; border-left: 2px solid #999; padding-left: 8px; }
        </style>
    """, unsafe_allow_html=True)


inject_custom_css()

# --- 2. 路径与配置 ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = os.path.join(BASE_DIR, 'data')
LOG_FILE = os.path.join(BASE_DIR, 'usage_log.csv')
CACHE_FILE = os.path.join(BASE_DIR, 'search_index.pkl')

USERS = {
    "admin": {"pwd": "admin123", "role": "admin", "dept": "管理中心", "name": "系统管理员"},
    "user1": {"pwd": "123456", "role": "user", "dept": "公共卫生团组", "name": "公卫专员"},
    "user2": {"pwd": "123456", "role": "user", "dept": "教育管理团组", "name": "教育专员"},
    "guest": {"pwd": "guest", "role": "guest", "dept": "访客", "name": "访客"},
}

SYNONYMS = {
    "登革热": ["伊蚊", "白纹伊蚊", "蚊媒传染病"],
    "高血压": ["慢性病", "心血管", "血压"],
}


# --- 3. 核心逻辑函数 ---

def init_log():
    cols = ['时间', '账号', '姓名', '部门', '操作类型', '对象/关键词', '详情']
    if not os.path.exists(LOG_FILE):
        pd.DataFrame(columns=cols).to_csv(LOG_FILE, index=False, encoding='utf-8-sig')
    else:
        try:
            pd.read_csv(LOG_FILE)
        except:
            pd.DataFrame(columns=cols).to_csv(LOG_FILE, index=False, encoding='utf-8-sig')


def log_action(action, target="", detail=""):
    try:
        init_log()
        u = st.session_state.get('username', 'Unknown')
        info = USERS.get(u, {})
        new_row = {
            '时间': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            '账号': u, '姓名': info.get('name', ''), '部门': info.get('dept', ''),
            '操作类型': action, '对象/关键词': target, '详情': detail
        }
        pd.DataFrame([new_row]).to_csv(LOG_FILE, mode='a', header=False, index=False, encoding='utf-8-sig')
    except:
        pass


# --- 索引构建 ---
@st.cache_resource(show_spinner=False)
def build_index():
    # 1. 尝试读取硬盘缓存
    if os.path.exists(CACHE_FILE):
        try:
            with open(CACHE_FILE, 'rb') as f:
                return pickle.load(f)
        except:
            pass

    index = []
    if not os.path.exists(DATA_DIR): return []

    # 扫描流程
    for root, dirs, files in os.walk(DATA_DIR):
        for file in files:
            ext = file.split('.')[-1].lower()
            if ext in ['pdf', 'docx', 'doc', 'xlsx', 'pptx']:
                full_path = os.path.join(root, file)
                path_parts = os.path.relpath(root, DATA_DIR).split(os.sep)

                # 提取团组
                dept_tag = "通用资源"
                for part in path_parts:
                    if "团组" in part: dept_tag = part; break

                year_match = re.search(r'202[0-9]', file)
                year_tag = year_match.group() if year_match else "----"

                # 读取内容
                content = ""
                try:
                    if ext == 'pdf':
                        with pdfplumber.open(full_path) as pdf:
                            if len(pdf.pages) > 0: content += pdf.pages[0].extract_text() or ""
                    elif ext in ['docx', 'doc']:
                        doc = docx.Document(full_path)
                        for p in doc.paragraphs[:20]: content += p.text + "\n"
                    elif ext == 'xlsx':
                        wb = openpyxl.load_workbook(full_path, data_only=True)
                        for sheet in wb.sheetnames:
                            for row in wb[sheet].iter_rows(max_row=5, values_only=True):
                                content += " ".join([str(c) for c in row if c]) + "\n"
                    elif ext == 'pptx':
                        prs = Presentation(full_path)
                        for slide in prs.slides[:3]:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"): content += shape.text + "\n"
                except:
                    content = ""

                index.append({
                    "name": file, "path": full_path, "type": ext,
                    "dept": dept_tag, "year": year_tag,
                    "category_str": " > ".join(path_parts),
                    "content": content
                })

    # 2. 保存到硬盘缓存
    try:
        with open(CACHE_FILE, 'wb') as f:
            pickle.dump(index, f)
    except:
        pass

    return index


# --- 辅助排序与展示 ---
def get_sorted_items(path):
    """文件/文件夹排序"""
    if not os.path.exists(path): return []
    items = [d for d in os.listdir(path) if not d.startswith('.')]
    items.sort()
    return items


def highlight_text(text, query):
    """关键词标黄处理"""
    if not query: return text
    pattern = re.compile(re.escape(query), re.IGNORECASE)
    return pattern.sub(f"<span style='background-color: #ffff00; color: #000; font-weight: bold;'>{query}</span>", text)


def render_file_row(file_name, file_path, context="browse", query=""):
    """渲染单行文件"""
    icon_map = {"pdf": "📕", "docx": "📘", "doc": "📘", "xlsx": "📗", "pptx": "📙"}
    ext = file_name.split('.')[-1].lower()
    icon = icon_map.get(ext, "📄")

    # 处理文件名显示
    display_name = file_name

    # 如果是搜索场景，处理关键词高亮
    if context == "search" and query:
        display_name = highlight_text(file_name, query)

    with st.container():
        c1, c2, c3 = st.columns([0.5, 8, 1.5])
        with c1: st.write(f"### {icon}")
        with c2: st.markdown(f"**{display_name}**", unsafe_allow_html=True)
        with c3:
            with open(file_path, "rb") as f:
                key = f"{context}_{file_path}_{int(time.time() * 10000)}"
                if st.download_button("📥 下载", f, file_name=file_name, key=key):
                    log_action("下载文件", file_name)
    st.markdown("<hr style='margin: 5px 0; border-top: 1px solid #eee;'>", unsafe_allow_html=True)


# --- 4. 页面功能模块 ---

def login_page():
    col1, col2, col3 = st.columns([1, 1, 1])
    with col2:
        st.markdown("<br><br>", unsafe_allow_html=True)
        with st.container(border=True):
            st.markdown("<h3 style='text-align:center; color:#1e50a2;'>系统用户登录</h3>", unsafe_allow_html=True)
            st.markdown("<hr>", unsafe_allow_html=True)
            u = st.text_input("用户名")
            p = st.text_input("密码", type="password")
            if st.button("登 录", use_container_width=True, type="primary"):
                if u in USERS and USERS[u]['pwd'] == p:
                    with st.spinner('正在验证身份...'):
                        st.session_state.update({'logged_in': True, 'username': u, 'role': USERS[u]['role']})
                        log_action("系统登录")
                        time.sleep(0.5)
                        st.rerun()
                else:
                    st.error("用户名或密码错误")
            st.markdown(
                "<div style='text-align:center; font-size:12px; color:#999; margin-top:20px;'>版权所有 © 2025 顺德区健康融入所有政策项目组</div>",
                unsafe_allow_html=True)


def main_app():
    # 顶部横幅 - 强制白色字体
    st.markdown(f"""
        <div style='background-color:#1e50a2; padding:15px; margin-bottom:20px; border-bottom: 3px solid #b22222;'>
            <div style='color:white !important; margin:0; font-size: 24px; font-weight: 600; letter-spacing: 2px;'>
                🏛️ 健康融入所有政策知识库平台
            </div>
        </div>
    """, unsafe_allow_html=True)

    # 侧边栏
    with st.sidebar:
        st.markdown(f"**当前用户：{USERS[st.session_state['username']]['name']}**")
        st.markdown("<hr>", unsafe_allow_html=True)
        # 名字：资源目录；检索查询；用户中心
        nav = st.radio("系统导航", ["资源目录", "检索查询", "用户中心"], label_visibility="collapsed")
        st.markdown("<hr>", unsafe_allow_html=True)

        if st.session_state['role'] == 'admin':
            if st.button("🔄 刷新数据库"):
                if os.path.exists(CACHE_FILE):
                    os.remove(CACHE_FILE)
                st.cache_resource.clear()
                st.rerun()
            st.markdown("<br>", unsafe_allow_html=True)

        if st.button("退出系统"):
            st.session_state['logged_in'] = False
            st.rerun()

    # --- 功能 1: 资源目录 (深度遍历支持第4级目录 - 纯文件夹展开形式) ---
    if nav == "资源目录":
        st.markdown("<div class='breadcrumb'>当前位置：首页 &gt; 资源目录</div>", unsafe_allow_html=True)

        if not os.path.exists(DATA_DIR):
            st.error("未找到 data 文件夹，请建立目录。")
            return

        # 获取一级目录 (Level 1)
        l1_dirs = [d for d in get_sorted_items(DATA_DIR) if os.path.isdir(os.path.join(DATA_DIR, d))]

        if not l1_dirs:
            st.info("知识库为空，请在 data 文件夹下建立目录。")

        # Level 1 使用 Tabs (保持不变)
        tabs = st.tabs(l1_dirs)

        for i, l1 in enumerate(l1_dirs):
            with tabs[i]:
                l1_path = os.path.join(DATA_DIR, l1)

                # 1. Level 1 直接文件
                l1_files = [f for f in get_sorted_items(l1_path) if
                            f.endswith(('.pdf', '.docx', '.doc', '.xlsx', '.pptx'))]
                if l1_files:
                    st.markdown("##### 📄 综合文档")
                    for f in l1_files: render_file_row(f, os.path.join(l1_path, f))

                # 2. Level 2 (二级文件夹) - 使用 Expander, 默认关闭
                l2_dirs = [d for d in get_sorted_items(l1_path) if os.path.isdir(os.path.join(l1_path, d))]

                if not l2_dirs and not l1_files:
                    st.caption("（此分类下暂无内容）")

                for l2 in l2_dirs:
                    with st.expander(f"📁 {l2}", expanded=False):  # 默认不展开
                        l2_path = os.path.join(l1_path, l2)

                        # 2.1 Level 2 直接文件
                        l2_files = [f for f in get_sorted_items(l2_path) if
                                    f.endswith(('.pdf', '.docx', '.doc', '.xlsx', '.pptx'))]
                        if l2_files:
                            for f in l2_files: render_file_row(f, os.path.join(l2_path, f))

                        # 2.2 Level 3 (三级子目录) - 使用 Expander, 默认关闭
                        l3_dirs = [d for d in get_sorted_items(l2_path) if os.path.isdir(os.path.join(l2_path, d))]

                        if not l3_dirs and not l2_files:
                            st.caption("（空文件夹）")

                        for l3 in l3_dirs:
                            with st.expander(f"🔹 {l3}", expanded=False):  # 嵌套 Expander, 默认不展开
                                l3_path = os.path.join(l2_path, l3)

                                # Level 3 直接文件
                                l3_files = [f for f in get_sorted_items(l3_path) if
                                            f.endswith(('.pdf', '.docx', '.doc', '.xlsx', '.pptx'))]

                                if l3_files:
                                    for f in l3_files: render_file_row(f, os.path.join(l3_path, f))

                                # 2.3 Level 4 (四级子目录) - 使用 Expander, 默认关闭
                                l4_dirs = [d for d in get_sorted_items(l3_path) if
                                           os.path.isdir(os.path.join(l3_path, d))]

                                if not l4_dirs and not l3_files:
                                    st.caption("（暂无文件）")

                                for l4 in l4_dirs:
                                    with st.expander(f"▪️ {l4}", expanded=False):  # 嵌套 Expander, 默认不展开
                                        l4_path = os.path.join(l3_path, l4)
                                        l4_files = [f for f in get_sorted_items(l4_path) if
                                                    f.endswith(('.pdf', '.docx', '.doc', '.xlsx', '.pptx'))]

                                        if l4_files:
                                            for f in l4_files: render_file_row(f, os.path.join(l4_path, f))
                                        else:
                                            st.markdown(
                                                "<div style='margin-left:20px;color:#999;font-size:12px'>（空文件夹）</div>",
                                                unsafe_allow_html=True)

                            st.markdown("<div style='margin-bottom:5px'></div>", unsafe_allow_html=True)

    # --- 功能 2: 检索查询 ---
    elif nav == "检索查询":
        st.markdown("<div class='breadcrumb'>当前位置：首页 &gt; 检索查询</div>", unsafe_allow_html=True)

        c1, c2 = st.columns([5, 1.5])
        query = c1.text_input("检索", placeholder="请输入关键词...", label_visibility="collapsed")
        mode = c2.radio("检索范围", ["仅标题", "搜全文"], horizontal=True, label_visibility="collapsed")
        st.caption("提示：'仅标题'只匹配文件名；'搜全文'匹配文件名及文档内容。")

        if 'db_index' not in st.session_state:
            st.session_state['db_index'] = build_index()
        index_data = st.session_state['db_index']

        if query:
            terms = [query]
            if query in SYNONYMS:
                terms.extend(SYNONYMS[query])
                st.info(f"💡 已启用智能联想：{', '.join(SYNONYMS[query])}")

            results = []
            for item in index_data:
                score = 0
                for term in terms:
                    term_lower = term.lower()
                    name_lower = item['name'].lower()
                    if mode == "仅标题":
                        if term_lower in name_lower: score += 10
                    else:
                        if term_lower in name_lower: score += 10
                        if term in item['content']: score += 5
                if score > 0:
                    results.append(item)

            st.markdown(f"**共检索到 {len(results)} 条记录**")
            st.markdown("<hr>", unsafe_allow_html=True)

            if results:
                for item in results:
                    render_file_row(item['name'], item['path'], context="search", query=query)
                    st.markdown(
                        f"<div style='color:#1e50a2; font-size:12px; margin-bottom:5px;'>📂 来源：{item['category_str']}</div>",
                        unsafe_allow_html=True)
            else:
                st.warning("暂无符合条件的数据")

    # --- 功能 3: 用户中心 ---
    elif nav == "用户中心":
        if st.session_state['role'] != 'admin':
            st.error("权限不足")
        else:
            st.markdown("<div class='breadcrumb'>当前位置：首页 &gt; 用户中心</div>", unsafe_allow_html=True)
            if os.path.exists(LOG_FILE):
                df = pd.read_csv(LOG_FILE)
                k1, k2, k3 = st.columns(3)
                k1.metric("总访问", f"{len(df)}")
                k2.metric("搜索量", f"{len(df[df['操作类型'].str.contains('检索', na=False)])}")
                k3.metric("下载量", f"{len(df[df['操作类型'].str.contains('下载', na=False)])}")

                st.markdown("##### 📊 部门活跃度")
                if not df.empty:
                    fig = px.pie(df, names='部门', hole=0.4)
                    fig.update_layout(margin=dict(t=0, b=0, l=0, r=0), height=300)
                    st.plotly_chart(fig, use_container_width=True)

                st.markdown("##### 📋 审计日志")
                st.dataframe(df.sort_index(ascending=False), use_container_width=True)


# --- 启动 ---
if __name__ == "__main__":
    if 'logged_in' not in st.session_state: st.session_state['logged_in'] = False
    if not st.session_state['logged_in']:
        login_page()
    else:
        main_app()